#!/usr/bin/env python3
"""头马即兴主持PPT - 12页完整版"""

import os
import zipfile
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# VI v1.3 色彩
WISDOM_BLUE = RGBColor(0x0D, 0x1A, 0x63)
BRIGHT_BLUE = RGBColor(0x06, 0x0E, 0x9F)
GOLD = RGBColor(0xFF, 0xB5, 0x47)
DARK_GREY = RGBColor(0x41, 0x40, 0x42)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)

def add_text(slide, left, top, width, height, text, font_size=Pt(16),
             font_color=DARK_GREY, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.name = 'Arial'
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    sp = shape._element
    style = sp.find(qn('p:style'))
    if style is not None:
        sp.remove(style)
    return shape

def add_hline(slide, x, y, length, color=WISDOM_BLUE, thickness=Pt(2)):
    return add_rect(slide, x, y, length, max(int(thickness), Emu(12700)), color)

def full_cleanup(outpath):
    tmppath = outpath + '.tmp'
    with zipfile.ZipFile(outpath, 'r') as zin:
        with zipfile.ZipFile(tmppath, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename.endswith('.xml'):
                    root = etree.fromstring(data)
                    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
                    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                    for style in root.findall(f'.//{{{ns_p}}}style'):
                        style.getparent().remove(style)
                    if 'theme' in item.filename.lower():
                        for tag in ['outerShdw', 'innerShdw', 'scene3d', 'sp3d']:
                            for el in root.findall(f'.//{{{ns_a}}}{tag}'):
                                el.getparent().remove(el)
                    data = etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone=True)
                zout.writestr(item, data)
    os.replace(tmppath, outpath)

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    
    LM = Inches(0.8)
    CW = Inches(11.733)
    
    # 第1页：流程介绍
    s1 = prs.slides.add_slide(blank)
    add_rect(s1, 0, 0, prs.slide_width, Inches(0.06), WISDOM_BLUE)
    add_text(s1, LM, Inches(0.8), CW, Inches(0.8), 
             "即兴演讲环节流程", font_size=Pt(36), font_color=WISDOM_BLUE, bold=True)
    add_hline(s1, LM, Inches(1.7), Inches(2), GOLD, Pt(3))
    
    steps = [
        ("01", "开场引言", "2分钟", "介绍'有力的发问'与今晚主题"),
        ("02", "第一轮·使命", "6分钟", "3位朋友 × 2分钟 | 你为什么存在？"),
        ("03", "第二轮·愿景", "6分钟", "3位朋友 × 2分钟 | 你要去向何方？"),
        ("04", "第三轮·价值观", "6分钟", "3位朋友 × 2分钟 | 你坚持什么？"),
        ("05", "收尾总结", "2分钟", "设计意图讲解 + 送给伙伴的话"),
    ]
    
    y_start = Inches(2.2)
    for i, (num, title, time, desc) in enumerate(steps):
        y = y_start + i * Inches(1.0)
        circle = s1.shapes.add_shape(MSO_SHAPE.OVAL, LM, y, Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = WISDOM_BLUE if i < 4 else GOLD
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        add_text(s1, LM + Inches(0.7), y, Inches(3), Inches(0.35),
                 title, font_size=Pt(18), font_color=WISDOM_BLUE, bold=True)
        add_text(s1, LM + Inches(3.8), y, Inches(1.2), Inches(0.3),
                 time, font_size=Pt(12), font_color=GOLD, bold=True)
        add_text(s1, LM + Inches(0.7), y + Inches(0.35), Inches(8), Inches(0.4),
                 desc, font_size=Pt(16), font_color=DARK_GREY)
    
    add_text(s1, LM, Inches(6.8), CW, Inches(0.4),
             "总时长：20分钟 | 参与人数：9位 | 柳州企业家头马第190次例会",
             font_size=Pt(12), font_color=DARK_GREY, alignment=PP_ALIGN.CENTER)
    
    # 第2页：前言
    s2 = prs.slides.add_slide(blank)
    add_rect(s2, 0, 0, prs.slide_width, Inches(0.06), WISDOM_BLUE)
    add_text(s2, LM, Inches(1.0), CW, Inches(0.8),
             "有力的发问", font_size=Pt(44), font_color=WISDOM_BLUE, bold=True)
    add_hline(s2, LM, Inches(1.9), Inches(1.5), GOLD, Pt(3))
    
    intro = """我做商业教练，有个核心工具叫"有力的发问"。

不是问"你怎么想的"，
是问一些让你停下来、重新看自己的问题。

今晚我想用这个工具，
带大家走一遍企业的"灵魂三问"。"""
    
    add_text(s2, LM, Inches(2.5), Inches(8), Inches(2.5),
             intro, font_size=Pt(16), font_color=DARK_GREY)
    
    cards = [("使命", "你为什么存在？"), ("愿景", "你要去向何方？"), ("价值观", "你坚持什么？")]
    for i, (t, st) in enumerate(cards):
        x = LM + i * Inches(4)
        add_rect(s2, x, Inches(5.0), Inches(3.5), Inches(0.06), BRIGHT_BLUE)
        add_text(s2, x, Inches(5.15), Inches(3.5), Inches(0.5),
                 t, font_size=Pt(22), font_color=BRIGHT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text(s2, x, Inches(5.7), Inches(3.5), Inches(0.4),
                 st, font_size=Pt(16), font_color=DARK_GREY, alignment=PP_ALIGN.CENTER)
    
    # 第3页：总议题（九宫格）
    s3 = prs.slides.add_slide(blank)
    add_rect(s3, 0, 0, prs.slide_width, Inches(0.06), WISDOM_BLUE)
    add_text(s3, LM, Inches(0.8), CW, Inches(0.8),
             "灵魂三问 · 九宫格", font_size=Pt(36), font_color=WISDOM_BLUE, bold=True)
    add_hline(s3, LM, Inches(1.7), Inches(2), GOLD, Pt(3))
    
    topics = [
        ("使命", ["第一个客户", "最骄傲的时刻", "如果重来"]),
        ("愿景", ["三年后的画面", "别人眼中的你", "想留下的东西"]),
        ("价值观", ["绝不妥协", "团队最像你的地方", "给新人的一句话"]),
    ]
    
    for row, (title, items) in enumerate(topics):
        y = Inches(2.2) + row * Inches(1.6)
        add_rect(s3, LM, y, Inches(2), Inches(0.06), BRIGHT_BLUE)
        add_text(s3, LM, y + Inches(0.15), Inches(2), Inches(0.5),
                 title, font_size=Pt(20), font_color=BRIGHT_BLUE, bold=True)
        
        for col, item in enumerate(items):
            x = LM + Inches(2.5) + col * Inches(3.5)
            add_rect(s3, x, y, Inches(3.3), Inches(1.3), LIGHT_BG)
            add_text(s3, x + Inches(0.2), y + Inches(0.2), Inches(3), Inches(0.9),
                     item, font_size=Pt(18), font_color=DARK_GREY, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 第4-9页：三轮议题详细页
    round_titles = [
        ("第一轮", "使命", "你为什么存在？", ["第一个客户", "最骄傲的时刻", "如果重来"]),
        ("第二轮", "愿景", "你要去向何方？", ["三年后的画面", "别人眼中的你", "想留下的东西"]),
        ("第三轮", "价值观", "你坚持什么？", ["绝不妥协", "团队最像你的地方", "给新人的一句话"]),
    ]
    
    for round_num, (rt, title, subtitle, items) in enumerate(round_titles):
        # 轮次标题页
        s_title = prs.slides.add_slide(blank)
        add_rect(s_title, 0, 0, prs.slide_width, Inches(0.06), WISDOM_BLUE)
        add_text(s_title, LM, Inches(2.0), CW, Inches(0.8),
                 f"{rt} · {title}", font_size=Pt(48), font_color=WISDOM_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text(s_title, LM, Inches(3.0), CW, Inches(0.6),
                 subtitle, font_size=Pt(28), font_color=BRIGHT_BLUE, alignment=PP_ALIGN.CENTER)
        add_hline(s_title, Inches(5), Inches(4.0), Inches(3.3), GOLD, Pt(3))
        
        # 关键词页
        s_detail = prs.slides.add_slide(blank)
        add_rect(s_detail, 0, 0, prs.slide_width, Inches(0.06), WISDOM_BLUE)
        add_text(s_detail, LM, Inches(0.8), CW, Inches(0.6),
                 f"{rt} · {title}", font_size=Pt(28), font_color=WISDOM_BLUE, bold=True)
        add_text(s_detail, LM, Inches(1.4), CW, Inches(0.4),
                 subtitle, font_size=Pt(18), font_color=BRIGHT_BLUE)
        add_hline(s_detail, LM, Inches(1.9), Inches(1.5), GOLD, Pt(2))
        
        for i, item in enumerate(items):
            y = Inches(2.5) + i * Inches(1.5)
            add_rect(s_detail, LM, y, Inches(0.8), Inches(0.8), BRIGHT_BLUE)
            num_box = s_detail.shapes.add_shape(MSO_SHAPE.OVAL, LM + Inches(0.1), y + Inches(0.1), Inches(0.6), Inches(0.6))
            num_box.fill.solid()
            num_box.fill.fore_color.rgb = WHITE
            num_box.line.fill.background()
            tf = num_box.text_frame
            tf.paragraphs[0].text = str(i + 1)
            tf.paragraphs[0].font.size = Pt(20)
            tf.paragraphs[0].font.color.rgb = BRIGHT_BLUE
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            add_text(s_detail, LM + Inches(1.2), y + Inches(0.15), Inches(8), Inches(0.6),
                     item, font_size=Pt(28), font_color=DARK_GREY, bold=True)
    
    # 第10页：设计意图
    s10 = prs.slides.add_slide(blank)
    add_rect(s10, 0, 0, prs.slide_width, Inches(0.06), WISDOM_BLUE)
    add_text(s10, LM, Inches(0.8), CW, Inches(0.8),
             "为什么这样设计？", font_size=Pt(36), font_color=WISDOM_BLUE, bold=True)
    add_hline(s10, LM, Inches(1.7), Inches(2), GOLD, Pt(3))
    
    design_text = """关键词而非问题 —— 给演讲者留白，让故事自然流淌

从使命到愿景到价值观 —— 循序渐进，由己及人

不点评只承上启下 —— 尊重每位演讲者，保持流程流畅

不销而销 —— 用"有力的发问"展示教练能力，而非推销产品"""
    
    add_text(s10, LM, Inches(2.3), Inches(10), Inches(3.5),
             design_text, font_size=Pt(18), font_color=DARK_GREY)
    
    # 第11页：送给伙伴的话
    s11 = prs.slides.add_slide(blank)
    add_rect(s11, 0, 0, prs.slide_width, Inches(0.06), WISDOM_BLUE)
    add_text(s11, LM, Inches(0.8), CW, Inches(0.8),
             "送给伙伴的话", font_size=Pt(36), font_color=WISDOM_BLUE, bold=True)
    add_hline(s11, LM, Inches(1.7), Inches(2), GOLD, Pt(3))
    
    quote = """九个人，九个答案。

没有标准答案，但每个答案都在说同一件事——
使命、愿景、价值观，不是写出来的，是选出来的。

今晚这些问题，我在教练对话里常问。
问对了，答案自己会出来。

如果某个问题让你想多聊几句，会后找我。"""
    
    add_text(s11, LM, Inches(2.3), Inches(10), Inches(3.5),
             quote, font_size=Pt(20), font_color=DARK_GREY)
    
    # 第12页：结尾
    s12 = prs.slides.add_slide(blank)
    add_rect(s12, 0, 0, prs.slide_width, Inches(0.06), WISDOM_BLUE)
    add_text(s12, LM, Inches(2.5), CW, Inches(1.0),
             "谢谢大家", font_size=Pt(56), font_color=WISDOM_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_hline(s12, Inches(5), Inches(3.6), Inches(3.3), GOLD, Pt(3))
    add_text(s12, LM, Inches(4.2), CW, Inches(0.6),
             "我是徐沛霖，商业教练", font_size=Pt(24), font_color=DARK_GREY, alignment=PP_ALIGN.CENTER)
    add_text(s12, LM, Inches(4.9), CW, Inches(0.5),
             "陪老板变强，让生意变好", font_size=Pt(18), font_color=BRIGHT_BLUE, alignment=PP_ALIGN.CENTER)
    
    # 保存
    outpath = '/Users/lucius/.openclaw/workspace/output/头马即兴主持_灵魂三问_PPT_v1.0.pptx'
    prs.save(outpath)
    full_cleanup(outpath)
    print(f'Created: {outpath} ({os.path.getsize(outpath):,} bytes)')

if __name__ == '__main__':
    main()
