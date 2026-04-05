#!/usr/bin/env python3
"""
小龙虾赋能营 · 宣传海报PPT
McKinsey-style Design System (v1.9.0)
"""

import os
import zipfile
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Color Palette ──
NAVY      = RGBColor(0x05, 0x1C, 0x2C)  # 深蓝 - 主色
BLACK     = RGBColor(0x00, 0x00, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY  = RGBColor(0x66, 0x66, 0x66)
LINE_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
BG_GRAY   = RGBColor(0xF2, 0xF2, 0xF2)

# Accent Colors
ACCENT_ORANGE = RGBColor(0xD4, 0x6A, 0x00)  # 橙 - 强调
ACCENT_BLUE   = RGBColor(0x00, 0x6B, 0xA6)  # 蓝
ACCENT_GREEN  = RGBColor(0x00, 0x7A, 0x53)  # 绿

# ── Font Sizes ──
TITLE_SIZE      = Pt(22)
BODY_SIZE       = Pt(14)
SUB_HEADER_SIZE = Pt(18)
HEADER_SIZE     = Pt(28)
SMALL_SIZE      = Pt(9)

# ── Helper Functions ──

def _clean_shape(shape):
    """Remove p:style from any shape to prevent effect references."""
    sp = shape._element
    style = sp.find(qn('p:style'))
    if style is not None:
        sp.remove(style)

def set_ea_font(run, typeface='KaiTi'):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', typeface)

def add_text(slide, left, top, width, height, text, font_size=Pt(14),
             font_name='Arial', font_color=DARK_GRAY, bold=False,
             alignment=PP_ALIGN.LEFT, ea_font='KaiTi', anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    anchor_map = {MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b', MSO_ANCHOR.TOP: 't'}
    bodyPr.set('anchor', anchor_map.get(anchor, 't'))
    for attr in ['lIns', 'tIns', 'rIns', 'bIns']:
        bodyPr.set(attr, '45720')
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.name = font_name
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    for run in p.runs:
        set_ea_font(run, ea_font)
    return txBox

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    _clean_shape(shape)
    return shape

def add_hline(slide, x, y, length, color=BLACK, thickness=Pt(0.5)):
    """Draw a horizontal line using a thin rectangle (no connector)."""
    h = max(int(thickness), Emu(6350))
    return add_rect(slide, x, y, length, h, color)

def add_oval(slide, x, y, letter, size=Inches(0.45), bg=NAVY, fg=WHITE):
    """Add a circle label with a letter."""
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    c.fill.solid()
    c.fill.fore_color.rgb = bg
    c.line.fill.background()
    tf = c.text_frame
    tf.paragraphs[0].text = letter
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.name = 'Arial'
    tf.paragraphs[0].font.color.rgb = fg
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for run in tf.paragraphs[0].runs:
        set_ea_font(run, 'KaiTi')
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    bodyPr.set('anchor', 'ctr')
    _clean_shape(c)
    return c

def add_action_title(slide, text, title_size=Pt(22)):
    add_text(slide, Inches(0.8), Inches(0.15), Inches(11.7), Inches(0.9),
             text, font_size=title_size, font_color=BLACK, bold=True,
             font_name='Georgia', ea_font='KaiTi', anchor=MSO_ANCHOR.MIDDLE)
    add_hline(slide, Inches(0.8), Inches(1.05), Inches(11.7),
              color=BLACK, thickness=Pt(0.5))

def add_source(slide, text, y=Inches(7.05)):
    add_text(slide, Inches(0.8), y, Inches(11), Inches(0.3),
             text, font_size=Pt(9), font_color=MED_GRAY)

def full_cleanup(outpath):
    """Remove ALL p:style from every slide + theme shadows/3D."""
    tmppath = outpath + '.tmp'
    with zipfile.ZipFile(outpath, 'r') as zin:
        with zipfile.ZipFile(tmppath, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename.endswith('.xml'):
                    root = etree.fromstring(data)
                    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
                    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                    for style in root.findall(f'.//{{{ns_p}}}style'):
                        style.getparent().remove(style)
                    if 'theme' in item.filename.lower():
                        for tag in ['outerShdw', 'innerShdw', 'scene3d', 'sp3d']:
                            for el in root.findall(f'.//{{{ns_a}}}{tag}'):
                                el.getparent().remove(el)
                    data = etree.tostring(root, xml_declaration=True,
                                          encoding='UTF-8', standalone=True)
                zout.writestr(item, data)
    os.replace(tmppath, outpath)

# ── Build Presentation ──

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1: Cover - 小龙虾赋能营主海报
    s1 = prs.slides.add_slide(blank)
    
    # 顶部装饰条
    add_rect(s1, 0, 0, prs.slide_width, Inches(0.08), NAVY)
    
    # 品牌标识区
    add_text(s1, Inches(0.8), Inches(0.3), Inches(4), Inches(0.4),
             '雨后甘霖', font_size=Pt(16), font_name='Georgia',
             font_color=NAVY, bold=True)
    add_text(s1, Inches(0.8), Inches(0.65), Inches(4), Inches(0.3),
             '阿徐商业教练', font_size=Pt(11), font_color=MED_GRAY)
    
    # 主标题
    add_text(s1, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.9),
             '小龙虾赋能营', font_size=Pt(48), font_name='Georgia',
             font_color=NAVY, bold=True)
    
    # 副标题
    add_text(s1, Inches(0.8), Inches(2.8), Inches(11.5), Inches(0.6),
             '4小时，让AI成为你的数字员工', font_size=Pt(28),
             font_color=DARK_GRAY)
    
    # 分隔线
    add_hline(s1, Inches(0.8), Inches(3.6), Inches(4), color=ACCENT_ORANGE, thickness=Pt(3))
    
    # 核心价值点
    add_text(s1, Inches(0.8), Inches(3.9), Inches(11.5), Inches(0.5),
             '省下2个员工年薪 · 掌握3大AI工具 · 带走50个即用模板', 
             font_size=Pt(16), font_color=MED_GRAY)
    
    # 课程信息卡片区域
    card_y = Inches(4.6)
    card_h = Inches(1.8)
    
    # 时间卡片
    add_rect(s1, Inches(0.8), card_y, Inches(3.5), card_h, BG_GRAY)
    add_rect(s1, Inches(0.8), card_y, Inches(3.5), Inches(0.06), ACCENT_BLUE)
    add_text(s1, Inches(1.0), card_y + Inches(0.15), Inches(3.1), Inches(0.35),
             '时间', font_size=Pt(12), font_color=MED_GRAY)
    add_text(s1, Inches(1.0), card_y + Inches(0.5), Inches(3.1), Inches(0.5),
             '2026.4.19', font_size=Pt(20), font_color=NAVY, bold=True)
    add_text(s1, Inches(1.0), card_y + Inches(1.0), Inches(3.1), Inches(0.4),
             '周六 14:00-18:00', font_size=Pt(13), font_color=DARK_GRAY)
    
    # 地点卡片
    add_rect(s1, Inches(4.6), card_y, Inches(3.5), card_h, BG_GRAY)
    add_rect(s1, Inches(4.6), card_y, Inches(3.5), Inches(0.06), ACCENT_GREEN)
    add_text(s1, Inches(4.8), card_y + Inches(0.15), Inches(3.1), Inches(0.35),
             '地点', font_size=Pt(12), font_color=MED_GRAY)
    add_text(s1, Inches(4.8), card_y + Inches(0.5), Inches(3.1), Inches(0.5),
             '柳州', font_size=Pt(20), font_color=NAVY, bold=True)
    add_text(s1, Inches(4.8), card_y + Inches(1.0), Inches(3.1), Inches(0.4),
             '具体地址报名后通知', font_size=Pt(13), font_color=DARK_GRAY)
    
    # 价格卡片
    add_rect(s1, Inches(8.4), card_y, Inches(4.0), card_h, NAVY)
    add_text(s1, Inches(8.6), card_y + Inches(0.15), Inches(3.6), Inches(0.35),
             '价格', font_size=Pt(12), font_color=RGBColor(0xAA, 0xAA, 0xAA))
    add_text(s1, Inches(8.6), card_y + Inches(0.45), Inches(3.6), Inches(0.5),
             '单人 ¥498', font_size=Pt(18), font_color=WHITE, bold=True)
    add_text(s1, Inches(8.6), card_y + Inches(0.95), Inches(3.6), Inches(0.4),
             '双人 ¥798', font_size=Pt(16), font_color=ACCENT_ORANGE, bold=True)
    
    # 限人数提示
    add_rect(s1, Inches(10.5), card_y + Inches(1.45), Inches(1.7), Inches(0.3), ACCENT_ORANGE)
    add_text(s1, Inches(10.5), card_y + Inches(1.45), Inches(1.7), Inches(0.3),
             '限30人', font_size=Pt(12), font_color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    
    # 底部CTA
    add_hline(s1, Inches(0.8), Inches(6.7), Inches(11.7), color=LINE_GRAY, thickness=Pt(0.5))
    add_text(s1, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.4),
             '扫码报名 · 先到先得 · 满员即止', 
             font_size=Pt(14), font_color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    # Save & cleanup
    outpath = '/Users/lucius/.openclaw/workspace/output/小龙虾赋能营_海报_v1.0.pptx'
    prs.save(outpath)
    full_cleanup(outpath)
    print(f'Created: {outpath} ({os.path.getsize(outpath):,} bytes)')

if __name__ == '__main__':
    main()