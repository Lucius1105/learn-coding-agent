#!/usr/bin/env python3
"""
Generate PPT for OpenClaw AI Enterprise Speech
简洁、高级、8页以内
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 宽屏
prs.slide_height = Inches(7.5)

# 配色方案 - Midnight Executive (专业克制)
COLORS = {
    'navy': RGBColor(30, 39, 97),        # 1E2761 主色
    'ice_blue': RGBColor(202, 220, 252), # CADCFC 辅色
    'white': RGBColor(255, 255, 255),    # FFFFFF
    'dark_gray': RGBColor(54, 69, 79),   # 深灰
    'light_gray': RGBColor(242, 242, 242), # 浅灰
    'gold': RGBColor(212, 175, 55),      # 金色点缀
}

def add_title_slide(prs, title, subtitle=""):
    """标题页 - 深色背景"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 深色背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['navy']
    background.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['ice_blue']
        p.alignment = PP_ALIGN.CENTER
    
    # 底部装饰线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4), Inches(6.5), Inches(5.333), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['gold']
    line.line.fill.background()
    
    return slide

def add_content_slide(prs, title, bullets, highlight_text="", layout="default"):
    """内容页 - 浅色背景"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 浅色背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['light_gray']
    background.line.fill.background()
    
    # 左侧装饰条
    sidebar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
        Inches(0.15), prs.slide_height
    )
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = COLORS['navy']
    sidebar.line.fill.background()
    
    # 页面标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['navy']
    
    # 高亮文字（大字）
    if highlight_text:
        hl_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(12), Inches(1.2))
        tf = hl_box.text_frame
        p = tf.paragraphs[0]
        p.text = highlight_text
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLORS['navy']
    
    # 要点列表
    if bullets:
        start_y = 3.2 if highlight_text else 1.8
        bullet_box = slide.shapes.add_textbox(Inches(0.8), Inches(start_y), Inches(12), Inches(4))
        tf = bullet_box.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(20)
            p.font.color.rgb = COLORS['dark_gray']
            p.space_after = Pt(16)
    
    return slide

def add_quote_slide(prs, quote, source=""):
    """金句页 - 深色背景"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 深色背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['navy']
    background.line.fill.background()
    
    # 引号装饰
    quote_mark = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(1), Inches(1))
    tf = quote_mark.text_frame
    p = tf.paragraphs[0]
    p.text = '"'
    p.font.size = Pt(120)
    p.font.color.rgb = COLORS['gold']
    
    # 金句内容
    quote_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.3), Inches(2))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = quote
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 来源
    if source:
        source_box = slide.shapes.add_textbox(Inches(1.5), Inches(5), Inches(10.3), Inches(0.6))
        tf = source_box.text_frame
        p = tf.paragraphs[0]
        p.text = source
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['ice_blue']
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_end_slide(prs, title, subtitle=""):
    """结尾页 - 深色背景"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 深色背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['navy']
    background.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.3), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.3), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(22)
        p.font.color.rgb = COLORS['ice_blue']
        p.alignment = PP_ALIGN.CENTER
    
    # 底部装饰线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4), Inches(6.5), Inches(5.333), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['gold']
    line.line.fill.background()
    
    return slide

# ============ 生成8页PPT ============

# 第1页：封面
add_title_slide(prs, 
    "用 OpenClaw 让企业跟上 AI 浪潮",
    "如何用 AI 实现短视频内容增量，解决获客难题")

# 第2页：开场演示
add_content_slide(prs,
    "这不是未来，这是现在",
    ["输入主题，10秒生成完整脚本", "开头钩子、核心内容、结尾金句", "拍摄建议、BGM推荐、封面设计"],
    "OpenClaw：AI 内容生产工具")

# 第3页：痛点
add_content_slide(prs,
    "AI 从搞钱开始",
    ["今天想选题，明天写文案，后天拍视频", "一条内容折腾两三天，发出去还没人看", "找人做，文案不像你，视频没灵魂"],
    "老板们想获客，但想到每天要生产内容，就一个头两个大")

# 第4页：核心观点
add_content_slide(prs,
    "内容获客的最大瓶颈",
    ["SOP：什么选题能获客，写成标准流程", "授权：团队做选题初稿，你只审脚本和出镜", "AI 化：用 OpenClaw 自动生成脚本，10秒出稿"],
    "不是创意，是流程")

# 第5页：三步流程
add_content_slide(prs,
    "SOP → 授权 → AI 化",
    ["梳理 SOP：目标客户画像、高转化选题、标准脚本结构", "明确授权：团队负责执行，你负责审核和出镜", "AI 化落地：OpenClaw 生成脚本，团队微调"],
    "你省下的时间，只做一件事：出镜")

# 第6页：案例
add_content_slide(prs,
    "真实案例：企业服务老板",
    ["之前：3个月8条，粉丝200，咨询0", "之后：每周5条，粉丝8000，每月20+咨询", "每天只花30分钟在内容上"],
    "从坚持不下来到持续产出")

# 第7页：金句
add_quote_slide(prs,
    "会用 AI 的老板，淘汰不会用的",
    "这个时代，获客的方式变了")

# 第8页：结尾
add_end_slide(prs,
    "欢迎来找我聊聊",
    "不是推销，只是聊聊。也许30分钟，能让你少走3个月的弯路")

# 保存文件
output_path = "/Users/lucius/.openclaw/workspace/output/演讲稿_OpenClaw_AI化企业_PPT_v1.0.pptx"
prs.save(output_path)
print(f"PPT已生成：{output_path}")
